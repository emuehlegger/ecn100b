#!/usr/bin/env python3
"""
Convert a lecture PPTX into Quarto markdown using the project's slide format rules.

Usage:
  py scripts/convert_lecture_slides_to_quarto.py input.pptx
  py scripts/convert_lecture_slides_to_quarto.py input.pptx -o s2026/lectures/lecture02.qmd
"""

from __future__ import annotations

import argparse
import re
from collections import defaultdict
from pathlib import Path


SLIDE_NUM_RE = re.compile(r"slide[_\-\s]*([0-9]+)", re.IGNORECASE)
ROMAN_SECTION_RE = re.compile(r"^\s*[IVXLC]+\.", re.IGNORECASE)


def normalize_text(text: str) -> str:
    return text.replace("\x0b", " ").strip()


def shape_position(shape) -> tuple[int, int]:
    return (getattr(shape, "top", 0), getattr(shape, "left", 0))


def get_slide_title(slide) -> str:
    title_shape = slide.shapes.title
    if title_shape is None:
        return "Untitled"
    title = normalize_text(title_shape.text or "")
    return title if title else "Untitled"


def extract_notes(slide) -> list[str]:
    notes_lines: list[str] = []
    if not slide.has_notes_slide:
        return notes_lines

    text_frame = slide.notes_slide.notes_text_frame
    if text_frame is None:
        return notes_lines

    for para in text_frame.paragraphs:
        txt = normalize_text(para.text or "")
        if txt:
            notes_lines.append(f"- {txt}")
    return notes_lines


def shape_text_lines(shape) -> list[str]:
    lines: list[str] = []
    if not shape.has_text_frame:
        return lines

    for para in shape.text_frame.paragraphs:
        txt = normalize_text(para.text or "")
        if not txt:
            continue

        level = max(0, int(getattr(para, "level", 0)))
        # Treat non-title text as bullets and preserve nesting depth.
        lines.append(f"{'  ' * level}- {txt}")
    return lines


def table_to_markdown(shape, slide_num: int) -> list[str]:
    if not shape.has_table:
        return []

    rows = []
    for row in shape.table.rows:
        cells = []
        for cell in row.cells:
            txt = normalize_text(cell.text or "")
            cells.append(txt.replace("|", "\\|"))
        rows.append(cells)

    if len(rows) < 2 or len(rows[0]) == 0:
        return [f"[TABLE: unable to extract cleanly; slide {slide_num}]"]

    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * len(rows[0])) + " |"
    body = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return [header, sep, *body]


def collect_svg_images_by_slide(images_dir: Path) -> dict[int, list[Path]]:
    slide_map: dict[int, list[Path]] = defaultdict(list)
    if not images_dir.exists():
        return slide_map

    for svg in images_dir.glob("*.svg"):
        match = SLIDE_NUM_RE.search(svg.name)
        if not match:
            continue
        slide_no = int(match.group(1))
        slide_map[slide_no].append(svg)

    for slide_no in slide_map:
        slide_map[slide_no].sort(key=lambda p: p.name.lower())
    return slide_map


def caption_from_filename(path: Path) -> str:
    stem = re.sub(r"[_\-]+", " ", path.stem).strip()
    return f"{stem} figure."


def alt_text_for_image(slide_no: int, slide_title: str, image_name: str) -> str:
    safe_title = slide_title.replace('"', "'")
    return (
        f"This image appears on slide {slide_no} titled '{safe_title}'. "
        f"It is sourced from {image_name} and is used as the corresponding figure for this slide. "
        "Describe axes, labels, arrows, and highlighted regions to explain the economic relationship shown."
    )


def front_matter(title: str) -> list[str]:
    return [
        "---",
        f'title: "{title}"',
        'subtitle: "Economics 100B - Intermediate Microeconomic Theory"',
        'author: "Prof. Muehlegger"',
        "format:",
        "  revealjs:",
        "    slide-level: 1",
        "---",
        "",
    ]


def build_slide_block(slide, slide_num: int, image_map: dict[int, list[Path]]) -> list[str]:
    title = get_slide_title(slide)

    body_lines: list[str] = []
    ordered_shapes = sorted(slide.shapes, key=shape_position)
    for shape in ordered_shapes:
        if shape == slide.shapes.title:
            continue

        if shape.has_table:
            body_lines.extend(table_to_markdown(shape, slide_num))
            body_lines.append("")
            continue

        text_lines = shape_text_lines(shape)
        if text_lines:
            body_lines.extend(text_lines)
            body_lines.append("")

    while body_lines and body_lines[-1] == "":
        body_lines.pop()

    heading_class = "{.medium-content}"
    if ROMAN_SECTION_RE.match(title) and not body_lines:
        heading_class = "{.center .medium-content}"

    out = [f"# {title} {heading_class}", ""]
    out.extend(body_lines)
    if body_lines:
        out.append("")

    for img in image_map.get(slide_num, []):
        caption = caption_from_filename(img)
        alt_text = alt_text_for_image(slide_num, title, img.name)
        out.append(
            f'![{caption}](./images/{img.name}){{fig-alt="{alt_text}"}}'
        )
        out.append("")

    notes = extract_notes(slide)
    if notes:
        out.append("Notes:")
        out.extend(notes)
        out.append("")

    while out and out[-1] == "":
        out.pop()
    out.append("")
    return out


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert lecture slides (.pptx) into a Quarto (.qmd) file."
    )
    parser.add_argument("pptx", type=Path, help="Path to input .pptx")
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=None,
        help="Output .qmd path (default: <pptx_basename>.qmd)",
    )
    parser.add_argument(
        "--images-dir",
        type=Path,
        default=None,
        help="Directory with extracted slide images (default: <output_dir>/images)",
    )
    args = parser.parse_args()

    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "Missing dependency 'python-pptx'. Install it with: py -m pip install python-pptx"
        ) from exc

    pptx_path = args.pptx.resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(f"Input file not found: {pptx_path}")
    if pptx_path.suffix.lower() != ".pptx":
        raise ValueError("Input file must be .pptx")

    output_path = args.output.resolve() if args.output else pptx_path.with_suffix(".qmd")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    images_dir = args.images_dir.resolve() if args.images_dir else output_path.parent / "images"
    image_map = collect_svg_images_by_slide(images_dir)

    pres = Presentation(str(pptx_path))
    title = pptx_path.stem

    lines = front_matter(title)
    for i, slide in enumerate(pres.slides, start=1):
        lines.extend(build_slide_block(slide, i, image_map))

    output_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")

    print(f"Created: {output_path}")
    print(f"Slides: {len(pres.slides)}")
    print(f"Images mapped (svg): {sum(len(v) for v in image_map.values())}")


if __name__ == "__main__":
    main()
