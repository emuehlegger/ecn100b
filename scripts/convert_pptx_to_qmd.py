#!/usr/bin/env python3
"""
convert_pptx_to_qmd.py

Convert a PowerPoint lecture deck to Quarto Markdown (.qmd) for Reveal.js.

Primary source : PPTX (slide order, text, math, notes) — authoritative.
Secondary source: PDF handout (math/encoding repair only).
Images         : ./Images folder with Lect{N}_Slide{M}* naming.

Usage:
    python scripts/convert_pptx_to_qmd.py "Class 1 - Introduction fv.pptx"
    python scripts/convert_pptx_to_qmd.py --pptx "file.pptx" [--pdf "handout.pdf"] [--images ./Images] [--out output.qmd]

Requirements:
    pip install python-pptx lxml
    pip install pymupdf          # optional, for PDF math repair
    # OR: pip install pdfminer.six  # fallback PDF backend
"""

from __future__ import annotations

import argparse
import os
import re
import sys
from pathlib import Path

# ── python-pptx ───────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from lxml import etree  # python-pptx already depends on lxml
except ImportError as exc:
    sys.exit(f"Missing dependency: {exc}\n  pip install python-pptx lxml")

# ── PDF backend (optional) ────────────────────────────────────────────────────
try:
    import fitz  # PyMuPDF
    _PDF_BACKEND = "pymupdf"
except ImportError:
    try:
        from pdfminer.high_level import extract_pages as _pdfminer_extract_pages
        from pdfminer.layout import LTTextContainer
        _PDF_BACKEND = "pdfminer"
    except ImportError:
        _PDF_BACKEND = None

# ── XML namespace constants ───────────────────────────────────────────────────
A_NS   = "http://schemas.openxmlformats.org/drawingml/2006/main"
MC_NS  = "http://schemas.openxmlformats.org/markup-compatibility/2006"
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
M_NS   = "http://schemas.openxmlformats.org/officeDocument/2006/math"
P_NS   = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _t(ns: str, local: str) -> str:
    """Return Clark-notation tag string."""
    return f"{{{ns}}}{local}"


# ── Unicode → LaTeX normalizer (applied to OMML output) ──────────────────────

# Greek letters
_GREEK_MAP: dict[str, str] = {
    "α": r"\alpha",   "β": r"\beta",    "γ": r"\gamma",   "δ": r"\delta",
    "ε": r"\epsilon", "ζ": r"\zeta",    "η": r"\eta",     "θ": r"\theta",
    "ι": r"\iota",    "κ": r"\kappa",   "λ": r"\lambda",  "μ": r"\mu",
    "ν": r"\nu",      "ξ": r"\xi",      "π": r"\pi",      "ρ": r"\rho",
    "σ": r"\sigma",   "τ": r"\tau",     "υ": r"\upsilon", "φ": r"\phi",
    "χ": r"\chi",     "ψ": r"\psi",     "ω": r"\omega",
    "Γ": r"\Gamma",   "Δ": r"\Delta",   "Θ": r"\Theta",   "Λ": r"\Lambda",
    "Ξ": r"\Xi",      "Π": r"\Pi",      "Σ": r"\Sigma",   "Υ": r"\Upsilon",
    "Φ": r"\Phi",     "Ψ": r"\Psi",     "Ω": r"\Omega",
    # Capital Greek that map to Roman letters in LaTeX (no backslash needed)
    "Α": "A", "Β": "B", "Ε": "E", "Ζ": "Z", "Η": "H",
    "Ι": "I", "Κ": "K", "Μ": "M", "Ν": "N", "Ο": "O",
    "Ρ": "R", "Τ": "T", "Χ": "X",
}

# Miscellaneous math symbols
_MATH_SYM_MAP: dict[str, str] = {
    "−": "-",               # U+2212 MINUS SIGN
    "×": r"\times",
    "÷": r"\div",
    "±": r"\pm",
    "∓": r"\mp",
    "≤": r"\leq",
    "≥": r"\geq",
    "≠": r"\neq",
    "≈": r"\approx",
    "≡": r"\equiv",
    "∈": r"\in",
    "∉": r"\notin",
    "⊂": r"\subset",
    "⊃": r"\supset",
    "⊆": r"\subseteq",
    "⊇": r"\supseteq",
    "∪": r"\cup",
    "∩": r"\cap",
    "∀": r"\forall",
    "∃": r"\exists",
    "∞": r"\infty",
    "→": r"\rightarrow",
    "←": r"\leftarrow",
    "⇒": r"\Rightarrow",
    "⇐": r"\Leftarrow",
    "⇔": r"\Leftrightarrow",
    "↔": r"\leftrightarrow",
    "↑": r"\uparrow",
    "↓": r"\downarrow",
    "∂": r"\partial",
    "∇": r"\nabla",
    "∗": "*",
    "·": r"\cdot",
    "…": r"\ldots",
    "⋯": r"\cdots",
    "⋮": r"\vdots",
    "⋱": r"\ddots",
    "′": r"'",
    "″": r"''",
    "√": r"\sqrt",
    "∝": r"\propto",
    "∼": r"\sim",
    "≃": r"\simeq",
    "∧": r"\wedge",
    "∨": r"\vee",
    "¬": r"\neg",
    "⊕": r"\oplus",
    "⊗": r"\otimes",
    "∅": r"\emptyset",
    "ℝ": r"\mathbb{R}",
    "ℤ": r"\mathbb{Z}",
    "ℕ": r"\mathbb{N}",
    "ℚ": r"\mathbb{Q}",
    "ℂ": r"\mathbb{C}",
}

# Regex to add backslash to bare math operator names
_MATH_OPS_RE = re.compile(
    r"(?<![\\a-zA-Z])(max|min|sup|inf|lim|arg|det|dim|exp|log|ln|"
    r"sin|cos|tan|arcsin|arccos|arctan|sinh|cosh|tanh|cot|sec|csc|"
    r"ker|rank|trace|tr|Pr|prob)(?![a-zA-Z])"
)


def _math_letter(cp: int) -> str | None:
    """Map a unicode code point to an ASCII letter if it's a math style variant."""
    # Math Italic Capital A-Z: U+1D434–U+1D44D
    if 0x1D434 <= cp <= 0x1D44D:
        return chr(ord("A") + cp - 0x1D434)
    # Math Italic Small a-z: U+1D44E–U+1D467 (hole at 0x1D455 for h)
    if 0x1D44E <= cp <= 0x1D467:
        off = cp - 0x1D44E
        return chr(ord("a") + off)
    # Math Bold Italic Capital: U+1D468–U+1D481
    if 0x1D468 <= cp <= 0x1D481:
        return chr(ord("A") + cp - 0x1D468)
    # Math Bold Italic Small: U+1D482–U+1D49B
    if 0x1D482 <= cp <= 0x1D49B:
        return chr(ord("a") + cp - 0x1D482)
    # Math Bold Capital: U+1D400–U+1D419
    if 0x1D400 <= cp <= 0x1D419:
        return chr(ord("A") + cp - 0x1D400)
    # Math Bold Small: U+1D41A–U+1D433
    if 0x1D41A <= cp <= 0x1D433:
        return chr(ord("a") + cp - 0x1D41A)
    # Math Sans-Serif Capital: U+1D5A0–U+1D5B9
    if 0x1D5A0 <= cp <= 0x1D5B9:
        return chr(ord("A") + cp - 0x1D5A0)
    # Math Sans-Serif Small: U+1D5BA–U+1D5D3
    if 0x1D5BA <= cp <= 0x1D5D3:
        return chr(ord("a") + cp - 0x1D5BA)
    # Math Sans-Serif Bold Capital: U+1D5D4–U+1D5ED
    if 0x1D5D4 <= cp <= 0x1D5ED:
        return chr(ord("A") + cp - 0x1D5D4)
    # Math Sans-Serif Bold Small: U+1D5EE–U+1D607
    if 0x1D5EE <= cp <= 0x1D607:
        return chr(ord("a") + cp - 0x1D5EE)
    # Math Script Capital: U+1D49C, U+1D49E, U+1D49F, U+1D4A2, U+1D4A5-6, U+1D4A9-AD, U+1D4AE-B7
    # (sparse — skip, handled by fallback)
    # PLANCK CONSTANT (h-bar context)
    if cp == 0x210E:
        return "h"
    return None


def _math_digit(cp: int) -> str | None:
    """Map a unicode code point to an ASCII digit if it's a math style variant."""
    for base, start in (
        (0x1D7CE, "0"),  # Math Bold
        (0x1D7D8, "0"),  # Math Double-Struck
        (0x1D7E2, "0"),  # Math Sans-Serif
        (0x1D7EC, "0"),  # Math Sans-Serif Bold
        (0x1D7F6, "0"),  # Math Monospace
    ):
        if base <= cp <= base + 9:
            return chr(ord("0") + cp - base)
    # Superscript digits
    superscripts = {0x00B2: "^{2}", 0x00B3: "^{3}", 0x00B9: "^{1}",
                    0x2070: "^{0}", 0x2074: "^{4}", 0x2075: "^{5}",
                    0x2076: "^{6}", 0x2077: "^{7}", 0x2078: "^{8}", 0x2079: "^{9}"}
    if cp in superscripts:
        return superscripts[cp]
    return None


def normalize_math_latex(s: str) -> str:
    """
    Post-process an OMML-derived LaTeX string:
    - Map unicode math-style letters to ASCII (already italic in math mode)
    - Map greek unicode to \\command, adding a trailing space when the next
      original character is alphanumeric (prevents \\PiR ambiguity)
    - Map math symbols to \\command, same trailing-space rule
    - Add \\ prefix to bare operator names (max, min, log, …)
    """
    result: list[str] = []
    chars = list(s)
    n = len(chars)
    for i, c in enumerate(chars):
        next_alnum = (i + 1 < n) and chars[i + 1].isalnum()

        if c in _GREEK_MAP:
            repl = _GREEK_MAP[c]
            result.append(repl)
            # If this becomes a \command and the next source char is a letter/digit,
            # add a space now so LaTeX won't merge the command with what follows.
            if repl.startswith("\\") and next_alnum:
                result.append(" ")
            continue

        if c in _MATH_SYM_MAP:
            repl = _MATH_SYM_MAP[c]
            result.append(repl)
            if repl.startswith("\\") and next_alnum:
                result.append(" ")
            continue

        cp = ord(c)
        letter = _math_letter(cp)
        if letter is not None:
            result.append(letter)
            continue

        digit = _math_digit(cp)
        if digit is not None:
            result.append(digit)
            continue

        result.append(c)

    normalized = "".join(result)
    # Add backslash to bare math operator names; then fix any adjacent-letter case
    # for the newly inserted backslashes (e.g. \maxq → \max q).
    normalized = _MATH_OPS_RE.sub(lambda m: "\\" + m.group(1), normalized)
    # Targeted space insertion only for known operator commands (avoids backtracking)
    normalized = re.sub(
        r"\\(max|min|sup|inf|lim|arg|det|dim|exp|log|ln"
        r"|sin|cos|tan|arcsin|arccos|arctan|sinh|cosh|tanh"
        r"|cot|sec|csc|ker|rank|trace|tr|Pr|prob)([A-Za-z])",
        r"\\\1 \2",
        normalized,
    )
    return normalized


# ── OMML → LaTeX converter ────────────────────────────────────────────────────

_NARY_CHR = {
    "∑": r"\sum",  "∏": r"\prod",
    "∫": r"\int",  "∬": r"\iint",  "∭": r"\iiint",
    "∮": r"\oint",
}
_ACC_CHR = {
    "^":  r"\hat",   "~":  r"\tilde", "→": r"\vec",
    "¯":  r"\bar",   "˙":  r"\dot",   "¨": r"\ddot",
    "\u20d7": r"\vec",
}
_BEG_DELIM = {
    "(": "(", "[": "[", "{": r"\{", "|": r"|",
    "⌈": r"\lceil", "⌊": r"\lfloor", "〈": r"\langle",
}
_END_DELIM = {
    ")": ")", "]": "]", "}": r"\}", "|": r"|",
    "⌉": r"\rceil", "⌋": r"\rfloor", "〉": r"\rangle",
}


def omml_to_latex(elem) -> str:
    """Recursively convert an OMML lxml element to a LaTeX string."""
    if elem is None:
        return ""
    tag = elem.tag

    def ch(local):
        return elem.find(_t(M_NS, local))

    def ch_all(local):
        return elem.findall(_t(M_NS, local))

    def sub_children():
        return "".join(omml_to_latex(c) for c in elem)

    # ── oMathPara / oMath ──
    if tag == _t(M_NS, "oMathPara"):
        maths = ch_all("oMath")
        return "\n".join(omml_to_latex(m) for m in maths) if maths else sub_children()

    if tag == _t(M_NS, "oMath"):
        return sub_children()

    # ── Math run ──
    if tag == _t(M_NS, "r"):
        t = ch("t")
        text = (t.text or "") if t is not None else ""
        # $ and % are special in LaTeX math; escape them if present literally.
        return text.replace("$", r"\$").replace("%", r"\%")

    # ── Fraction ──
    if tag == _t(M_NS, "f"):
        num = omml_to_latex(ch("num"))
        den = omml_to_latex(ch("den"))
        return rf"\frac{{{num}}}{{{den}}}"

    # ── Radical ──
    if tag == _t(M_NS, "rad"):
        radPr = ch("radPr")
        hide_deg = False
        if radPr is not None:
            dh = radPr.find(_t(M_NS, "degHide"))
            if dh is not None:
                hide_deg = dh.get(_t(M_NS, "val"), "0") not in ("0", "false")
        content = omml_to_latex(ch("e"))
        if hide_deg:
            return rf"\sqrt{{{content}}}"
        deg_e = ch("deg")
        deg_str = omml_to_latex(deg_e).strip() if deg_e is not None else ""
        if not deg_str or deg_str == "2":
            return rf"\sqrt{{{content}}}"
        return rf"\sqrt[{deg_str}]{{{content}}}"

    # ── Superscript ──
    if tag == _t(M_NS, "sSup"):
        base = omml_to_latex(ch("e"))
        sup  = omml_to_latex(ch("sup"))
        b = f"{{{base}}}" if len(base) > 1 else base
        return f"{b}^{{{sup}}}"

    # ── Subscript ──
    if tag == _t(M_NS, "sSub"):
        base = omml_to_latex(ch("e"))
        sub  = omml_to_latex(ch("sub"))
        b = f"{{{base}}}" if len(base) > 1 else base
        return f"{b}_{{{sub}}}"

    # ── Sub-superscript ──
    if tag == _t(M_NS, "sSubSup"):
        base = omml_to_latex(ch("e"))
        sub  = omml_to_latex(ch("sub"))
        sup  = omml_to_latex(ch("sup"))
        return f"{base}_{{{sub}}}^{{{sup}}}"

    # ── N-ary operator (∑, ∫, …) ──
    if tag == _t(M_NS, "nary"):
        nPr = ch("naryPr")
        chr_e = nPr.find(_t(M_NS, "chr")) if nPr is not None else None
        chr_val = chr_e.get(_t(M_NS, "val"), "∑") if chr_e is not None else "∑"
        op = _NARY_CHR.get(chr_val, chr_val)
        sub_e = ch("sub"); sup_e = ch("sup"); e = ch("e")
        sub_str = f"_{{{omml_to_latex(sub_e)}}}" if sub_e is not None else ""
        sup_str = f"^{{{omml_to_latex(sup_e)}}}" if sup_e is not None else ""
        content = omml_to_latex(e) if e is not None else ""
        return f"{op}{sub_str}{sup_str}{content}"

    # ── Function (e.g. max, min) ──
    if tag == _t(M_NS, "func"):
        fname = omml_to_latex(ch("fName"))
        arg   = omml_to_latex(ch("e"))
        return f"{fname}({arg})"

    # ── Delimiter ──
    if tag == _t(M_NS, "d"):
        dPr = ch("dPr")
        if dPr is not None:
            bc = dPr.find(_t(M_NS, "begChr"))
            ec = dPr.find(_t(M_NS, "endChr"))
            beg = bc.get(_t(M_NS, "val"), "(") if bc is not None else "("
            end = ec.get(_t(M_NS, "val"), ")") if ec is not None else ")"
        else:
            beg, end = "(", ")"
        beg_l = _BEG_DELIM.get(beg, beg)
        end_l = _END_DELIM.get(end, end)
        parts = [omml_to_latex(e) for e in ch_all("e")]
        return rf"\left{beg_l}" + ", ".join(parts) + rf"\right{end_l}"

    # ── Equation array ──
    if tag == _t(M_NS, "eqArr"):
        rows = [omml_to_latex(e) for e in ch_all("e")]
        return "\\begin{aligned}\n" + " \\\\\n".join(rows) + "\n\\end{aligned}"

    # ── Matrix ──
    if tag == _t(M_NS, "m"):
        rows = []
        for mr in elem.findall(_t(M_NS, "mr")):
            cells = [omml_to_latex(e) for e in mr.findall(_t(M_NS, "e"))]
            rows.append(" & ".join(cells))
        return "\\begin{pmatrix}\n" + " \\\\\n".join(rows) + "\n\\end{pmatrix}"

    # ── Limits ──
    if tag == _t(M_NS, "limLow"):
        return omml_to_latex(ch("e")) + f"_{{{omml_to_latex(ch('lim'))}}}"

    if tag == _t(M_NS, "limUpp"):
        return omml_to_latex(ch("e")) + f"^{{{omml_to_latex(ch('lim'))}}}"

    # ── Accent (hat, tilde, vec …) ──
    if tag == _t(M_NS, "acc"):
        accPr = ch("accPr")
        chr_e = accPr.find(_t(M_NS, "chr")) if accPr is not None else None
        chr_val = chr_e.get(_t(M_NS, "val"), "^") if chr_e is not None else "^"
        cmd = _ACC_CHR.get(chr_val, r"\hat")
        return f"{cmd}{{{omml_to_latex(ch('e'))}}}"

    # ── Bar / overline ──
    if tag == _t(M_NS, "bar"):
        barPr = ch("barPr")
        pos_e = barPr.find(_t(M_NS, "pos")) if barPr is not None else None
        pos = pos_e.get(_t(M_NS, "val"), "top") if pos_e is not None else "top"
        inner = omml_to_latex(ch("e"))
        return (rf"\overline{{{inner}}}" if pos == "top" else rf"\underline{{{inner}}}")

    # ── Box / group-char (pass-through) ──
    if tag in (_t(M_NS, "box"), _t(M_NS, "groupChr")):
        return omml_to_latex(ch("e"))

    # ── fName / num / den / e / sub / sup … (structural wrappers) ──
    if tag in (
        _t(M_NS, "fName"), _t(M_NS, "num"), _t(M_NS, "den"),
        _t(M_NS, "e"), _t(M_NS, "sub"), _t(M_NS, "sup"),
        _t(M_NS, "deg"), _t(M_NS, "lim"),
    ):
        return sub_children()

    # ── Fallback: recurse ──
    return sub_children()


# ── PDF text extraction ───────────────────────────────────────────────────────

def extract_pdf_pages(pdf_path: Path) -> list[str]:
    """Return list of text strings, one per page. Empty list if unavailable."""
    if _PDF_BACKEND is None or not pdf_path or not pdf_path.exists():
        return []
    if _PDF_BACKEND == "pymupdf":
        doc = fitz.open(str(pdf_path))
        pages = [page.get_text("text") for page in doc]
        doc.close()
        return pages
    else:  # pdfminer
        pages = []
        for page_layout in _pdfminer_extract_pages(str(pdf_path)):
            text = ""
            for el in page_layout:
                if isinstance(el, LTTextContainer):
                    text += el.get_text()
            pages.append(text)
        return pages


# ── Paragraph token extraction ────────────────────────────────────────────────

class _Token:
    __slots__ = ("kind", "content")

    def __init__(self, kind: str, content: str):
        self.kind = kind        # "text" | "math_inline" | "math_display"
        self.content = content


def _para_indent_level(p_elem) -> int:
    pPr = p_elem.find(_t(A_NS, "pPr"))
    if pPr is None:
        return 0
    try:
        return int(pPr.get("lvl", "0"))
    except ValueError:
        return 0


def _para_is_bulleted(p_elem) -> bool:
    """Return True if the paragraph explicitly carries a bullet marker."""
    pPr = p_elem.find(_t(A_NS, "pPr"))
    if pPr is None:
        return False
    if pPr.find(_t(A_NS, "buNone")) is not None:
        return False
    if pPr.find(_t(A_NS, "buChar")) is not None:
        return True
    if pPr.find(_t(A_NS, "buAutoNum")) is not None:
        return True
    if pPr.find(_t(A_NS, "buFont")) is not None:
        return True
    return False


def _para_autonumber_type(p_elem) -> str | None:
    """Return the buAutoNum type string if this paragraph uses auto-numbering, else None."""
    pPr = p_elem.find(_t(A_NS, "pPr"))
    if pPr is None:
        return None
    buAutoNum = pPr.find(_t(A_NS, "buAutoNum"))
    if buAutoNum is None:
        return None
    return buAutoNum.get("type", "arabicPeriod")


def _int_to_alpha(n: int) -> str:
    """1→a, 2→b, …, 26→z, 27→aa, etc."""
    result = ""
    while n > 0:
        n -= 1
        result = chr(ord("a") + n % 26) + result
        n //= 26
    return result


def _int_to_roman(n: int) -> str:
    vals = [
        (1000, "m"), (900, "cm"), (500, "d"), (400, "cd"),
        (100, "c"),  (90, "xc"),  (50, "l"),  (40, "xl"),
        (10, "x"),   (9, "ix"),   (5, "v"),   (4, "iv"),  (1, "i"),
    ]
    result = ""
    for v, s in vals:
        while n >= v:
            result += s
            n -= v
    return result


def _make_list_label(autonumber_type: str, count: int) -> str:
    """Return the Markdown list prefix (e.g. 'A.' or 'iii)') for the nth item."""
    _TYPE_MAP = {
        "arabicPeriod":  (str,           False, "."),
        "arabicParenR":  (str,           False, ")"),
        "alphaLcPeriod": (_int_to_alpha, False, "."),
        "alphaUcPeriod": (_int_to_alpha, True,  "."),
        "alphaLcParenR": (_int_to_alpha, False, ")"),
        "alphaUcParenR": (_int_to_alpha, True,  ")"),
        "romanLcPeriod": (_int_to_roman, False, "."),
        "romanUcPeriod": (_int_to_roman, True,  "."),
        "romanLcParenR": (_int_to_roman, False, ")"),
        "romanUcParenR": (_int_to_roman, True,  ")"),
    }
    if autonumber_type in _TYPE_MAP:
        fn, upper, sep = _TYPE_MAP[autonumber_type]
        if fn is str:
            label = str(count)
        else:
            label = fn(count)
            if upper:
                label = label.upper()
        return f"{label}{sep}"
    return f"{count}."  # fallback


def _extract_para_tokens(p_elem) -> list[_Token]:
    """
    Walk the children of <a:p> and return a list of _Token.
    Handles interleaved text runs and OMML math blocks.
    """
    tokens: list[_Token] = []
    cur_text = ""

    def flush():
        nonlocal cur_text
        s = cur_text
        cur_text = ""
        if s:
            tokens.append(_Token("text", s))

    for child in p_elem:
        ctag = child.tag

        # Regular text run
        if ctag == _t(A_NS, "r"):
            t = child.find(_t(A_NS, "t"))
            if t is not None and t.text:
                cur_text += t.text

        # Direct a14:m math element (common when slide content placeholder
        # is wrapped in mc:AlternateContent at the spTree level — python-pptx
        # then exposes paragraphs where math sits directly as <a14:m>)
        elif ctag == _t(A14_NS, "m"):
            flush()
            oMathPara = child.find(_t(M_NS, "oMathPara"))
            if oMathPara is not None:
                oPr = oMathPara.find(_t(M_NS, "oMathParaPr"))
                jc_val = ""
                if oPr is not None:
                    jc_e = oPr.find(_t(M_NS, "jc"))
                    if jc_e is not None:
                        jc_val = jc_e.get(_t(M_NS, "val"), "")
                latex = normalize_math_latex(omml_to_latex(oMathPara).strip())
                kind = "math_inline" if jc_val == "inline" else "math_display"
                tokens.append(_Token(kind, latex))
            else:
                oMath = child.find(_t(M_NS, "oMath"))
                if oMath is not None:
                    latex = normalize_math_latex(omml_to_latex(oMath).strip())
                    tokens.append(_Token("math_inline", latex))

        # AlternateContent — may contain OMML math
        elif ctag == _t(MC_NS, "AlternateContent"):
            choice = child.find(_t(MC_NS, "Choice"))
            math_found = False
            if choice is not None:
                a14m = choice.find(_t(A14_NS, "m"))
                if a14m is not None:
                    flush()
                    math_found = True
                    # Determine inline vs display from oMathParaPr/jc
                    oMathPara = a14m.find(_t(M_NS, "oMathPara"))
                    if oMathPara is not None:
                        oPr = oMathPara.find(_t(M_NS, "oMathParaPr"))
                        jc_val = ""
                        if oPr is not None:
                            jc_e = oPr.find(_t(M_NS, "jc"))
                            if jc_e is not None:
                                jc_val = jc_e.get(_t(M_NS, "val"), "")
                        latex = normalize_math_latex(omml_to_latex(oMathPara).strip())
                        # "inline" → inline; anything else (center/left/right/empty) → display
                        kind = "math_inline" if jc_val == "inline" else "math_display"
                        tokens.append(_Token(kind, latex))
                    else:
                        oMath = a14m.find(_t(M_NS, "oMath"))
                        if oMath is not None:
                            latex = normalize_math_latex(omml_to_latex(oMath).strip())
                            tokens.append(_Token("math_inline", latex))

            if not math_found:
                # Fallback: extract plain text from mc:Fallback
                fb = child.find(_t(MC_NS, "Fallback"))
                if fb is not None:
                    for r in fb.findall(f".//{_t(A_NS, 'r')}"):
                        t = r.find(_t(A_NS, "t"))
                        if t is not None and t.text:
                            cur_text += t.text

        # Field (slide number, date, etc.)
        elif ctag == _t(A_NS, "fld"):
            t = child.find(_t(A_NS, "t"))
            if t is not None and t.text:
                cur_text += t.text

        # Line break
        elif ctag == _t(A_NS, "br"):
            cur_text += "\n"

        # pPr, endParaMark — skip
    flush()
    return tokens


def _tokens_to_md(tokens: list[_Token]) -> str:
    """Convert token list to a markdown string (for a single paragraph)."""
    parts = []
    for tok in tokens:
        if tok.kind == "text":
            # Escape $ and % so they are not misread as LaTeX math delimiters
            # or comment characters in the rendered Markdown/Quarto output.
            text = tok.content.replace("$", r"\$").replace("%", r"\%")
            parts.append(text)
        elif tok.kind == "math_inline":
            parts.append(f"${tok.content}$")
        elif tok.kind == "math_display":
            # Will be emitted as display on its own line later;
            # mark with sentinel so caller can split
            parts.append(f"\x00DISPLAY\x00{tok.content}\x00END\x00")
    return "".join(parts)


# ── XML-based shape helpers ───────────────────────────────────────────────────
# Maps XML ph/@type attribute → integer code (matching PP_PLACEHOLDER enum)
_PH_TYPE_STR_TO_INT: dict[str, int] = {
    "title": 1, "ctrTitle": 3, "subTitle": 4, "body": 2,
    "dt": 16, "ftr": 15, "sldNum": 13, "obj": 7,
    "chart": 8, "dgm": 11, "media": 12, "clipArt": 10, "pic": 18,
}
_CHROME_PH_TYPES = frozenset({1, 3, 13, 15, 16})  # title, ctrTitle, sldNum, ftr, dt
_TABLE_GRAPHIC_URI = "http://schemas.openxmlformats.org/drawingml/2006/table"


def sp_elem_ph_type(sp_elem) -> int | None:
    """Get placeholder type from a p:sp XML element, or None if not a placeholder."""
    nvPr = sp_elem.find(f".//{_t(P_NS, 'nvPr')}")
    if nvPr is None:
        return None
    ph = nvPr.find(_t(P_NS, "ph"))
    if ph is None:
        return None
    type_str = ph.get("type", "")
    if type_str:
        return _PH_TYPE_STR_TO_INT.get(type_str, 0)
    # ph element present but no type → index-based content placeholder (body)
    return 2


def sp_elem_position(elem) -> tuple[int, int]:
    """Return (top, left) in EMU for sort ordering."""
    xfrm = elem.find(f".//{_t(A_NS, 'xfrm')}")
    if xfrm is None:
        return (0, 0)
    off = xfrm.find(_t(A_NS, "off"))
    if off is None:
        return (0, 0)
    return (int(off.get("y", 0)), int(off.get("x", 0)))


def _graphicframe_is_table(gf_elem) -> bool:
    """Return True if a graphicFrame element contains a table."""
    gd = gf_elem.find(f".//{_t(A_NS, 'graphicData')}")
    if gd is None:
        return False
    return gd.get("uri", "") == _TABLE_GRAPHIC_URI


def _iter_sp_elems_in(container):
    """Generator: yield (tag, elem, ph_type, top, left) for shapes in container."""
    for child in container:
        tag = child.tag
        if tag == _t(P_NS, "sp"):
            ph_type = sp_elem_ph_type(child)
            top, left = sp_elem_position(child)
            yield (tag, child, ph_type, top, left)
        elif tag in (_t(P_NS, "pic"), _t(P_NS, "graphicFrame")):
            top, left = sp_elem_position(child)
            yield (tag, child, None, top, left)
        elif tag == _t(MC_NS, "AlternateContent"):
            # Unwrap mc:Choice — python-pptx often misses these
            choice = child.find(_t(MC_NS, "Choice"))
            if choice is not None:
                yield from _iter_sp_elems_in(choice)
        elif tag == _t(P_NS, "grpSp"):
            yield from _iter_sp_elems_in(child)


def iter_all_shape_elems(slide):
    """
    Yield (tag, elem, ph_type, top, left) for every shape in the slide's spTree,
    including shapes wrapped in mc:AlternateContent that python-pptx would miss.
    """
    spTree = slide._element.find(f"{_t(P_NS, 'cSld')}/{_t(P_NS, 'spTree')}")
    if spTree is None:
        return
    yield from _iter_sp_elems_in(spTree)


# ── Shape content extraction ──────────────────────────────────────────────────

class SlideItem:
    """A single piece of content for a slide."""
    __slots__ = ("kind", "content", "level", "is_bullet", "list_label")

    def __init__(
        self,
        kind: str,
        content: str,
        level: int = 0,
        is_bullet: bool = False,
        list_label: str | None = None,
    ):
        self.kind = kind   # "text" | "math_display" | "empty"
        self.content = content
        self.level = level
        self.is_bullet = is_bullet
        self.list_label = list_label  # e.g. "A." or "1)" for autonumbered lists


_DISPLAY_RE = re.compile(r"\x00DISPLAY\x00(.*?)\x00END\x00", re.DOTALL)


def extract_text_from_txbody(txBody, is_body: bool) -> list[SlideItem]:
    """Extract SlideItems from a p:txBody lxml element."""
    items: list[SlideItem] = []
    # Counters for buAutoNum sequences, keyed by (level, type_str).
    autonumber_counters: dict[tuple[int, str], int] = {}

    for p_elem in txBody.findall(_t(A_NS, "p")):
        level = _para_indent_level(p_elem)
        explicit_bullet = _para_is_bulleted(p_elem)
        autonumber_type = _para_autonumber_type(p_elem)

        pPr = p_elem.find(_t(A_NS, "pPr"))
        has_buNone = pPr is not None and pPr.find(_t(A_NS, "buNone")) is not None
        is_bullet = (
            explicit_bullet
            or (is_body and level == 0 and not has_buNone)
            or level > 0
        )

        tokens = _extract_para_tokens(p_elem)
        if not tokens:
            items.append(SlideItem("empty", "", level, False))
            continue

        md = _tokens_to_md(tokens)
        segments = _DISPLAY_RE.split(md)

        # Determine list_label for the first text segment (autonumber items only).
        # Only increment the counter for paragraphs that have actual text content.
        list_label: str | None = None
        if autonumber_type:
            has_text = any(seg.strip() for idx, seg in enumerate(segments) if idx % 2 == 0)
            if has_text:
                key = (level, autonumber_type)
                autonumber_counters[key] = autonumber_counters.get(key, 0) + 1
                list_label = _make_list_label(autonumber_type, autonumber_counters[key])

        for idx, seg in enumerate(segments):
            if idx % 2 == 1:
                items.append(SlideItem("math_display", seg.strip(), level, False))
            else:
                text = seg.strip()
                if text:
                    items.append(SlideItem(
                        "text", text, level,
                        is_bullet and idx == 0,
                        list_label if idx == 0 else None,
                    ))

    return items


def extract_table_from_elem(gf_elem, slide_num: int) -> str:
    """Extract Markdown table from a graphicFrame XML element."""
    try:
        tbl = gf_elem.find(f".//{_t(A_NS, 'tbl')}")
        if tbl is None:
            return f"[TABLE: unable to extract cleanly; slide {slide_num}]"
        rows = []
        for tr in tbl.findall(_t(A_NS, "tr")):
            cells = []
            for tc in tr.findall(_t(A_NS, "tc")):
                cell_text = ""
                txBody = tc.find(_t(A_NS, "txBody"))
                if txBody is not None:
                    parts = []
                    for p in txBody.findall(_t(A_NS, "p")):
                        line_tokens = _extract_para_tokens(p)
                        parts.append(_tokens_to_md(line_tokens))
                    cell_text = " ".join(p.strip() for p in parts if p.strip())
                cells.append(cell_text)
            rows.append(cells)
        if not rows:
            return f"[TABLE: empty; slide {slide_num}]"
        header = rows[0]
        md_rows = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * len(header)) + " |",
        ]
        for row in rows[1:]:
            md_rows.append("| " + " | ".join(row) + " |")
        return "\n".join(md_rows)
    except Exception:
        return f"[TABLE: unable to extract cleanly; slide {slide_num}]"


# ── Image mapping ─────────────────────────────────────────────────────────────

def find_slide_images(images_dir: Path, lect_num: int, slide_num: int) -> list[Path]:
    """
    Find image files for lecture lect_num, slide slide_num.
    Returns one file per unique image, preferring SVG > PNG > JPG.
    """
    if not images_dir or not images_dir.exists():
        return []

    pattern = re.compile(
        rf"[Ll]ect{re.escape(str(lect_num))}_[Ss]lide{re.escape(str(slide_num))}(?:\D|$)",
    )

    ext_priority = {".svg": 0, ".png": 1, ".jpg": 2, ".jpeg": 3}
    best: dict[str, tuple[int, Path]] = {}  # normalized_stem → (priority, path)

    for f in images_dir.iterdir():
        if not f.is_file():
            continue
        if not pattern.search(f.stem):
            continue
        ext = f.suffix.lower()
        if ext not in ext_priority:
            continue
        # Normalize: strip extension + lowercase for case-insensitive dedup
        norm = re.sub(r"\.(svg|png|jpg|jpeg)$", "", f.name, flags=re.IGNORECASE).lower()
        prio = ext_priority[ext]
        if norm not in best or prio < best[norm][0]:
            best[norm] = (prio, f)

    return [v[1] for v in sorted(best.values(), key=lambda x: x[1].name)]


# ── Section-divider detection ─────────────────────────────────────────────────

_SECTION_RE = re.compile(r"^[IVXLCDM]+\.\s", re.IGNORECASE)


def is_section_divider(title: str) -> bool:
    return bool(_SECTION_RE.match(title.strip()))


# ── Slide title / notes ───────────────────────────────────────────────────────

def _normalize_title(text: str) -> str:
    """Collapse runs of whitespace (including newlines/tabs) to single spaces."""
    return re.sub(r"\s+", " ", text).strip()


def get_slide_title(slide) -> str:
    """Extract slide title, checking all shapes including mc:AlternateContent-wrapped ones."""
    # 1. Prefer explicitly-typed title placeholders
    for tag, elem, ph_type, top, left in iter_all_shape_elems(slide):
        if tag != _t(P_NS, "sp") or ph_type not in (1, 3):
            continue
        txBody = elem.find(_t(P_NS, "txBody"))
        if txBody is None:
            continue
        texts = []
        for p in txBody.findall(_t(A_NS, "p")):
            toks = _extract_para_tokens(p)
            texts.append(_tokens_to_md(toks))
        t = _normalize_title(" ".join(texts))
        if t:
            return t
    # 2. Fallback: topmost non-chrome text shape
    candidates = []
    for tag, elem, ph_type, top, left in iter_all_shape_elems(slide):
        if tag != _t(P_NS, "sp") or ph_type in _CHROME_PH_TYPES:
            continue
        txBody = elem.find(_t(P_NS, "txBody"))
        if txBody is None:
            continue
        texts = []
        for p in txBody.findall(_t(A_NS, "p")):
            toks = _extract_para_tokens(p)
            texts.append(_tokens_to_md(toks))
        t = _normalize_title(" ".join(texts))
        if t:
            candidates.append((top, left, t))
    if candidates:
        return sorted(candidates)[0][2]
    return ""


def get_slide_notes(slide) -> list[str]:
    try:
        if not slide.has_notes_slide:
            return []
        tf = slide.notes_slide.notes_text_frame
        if tf is None:
            return []
        lines = []
        for para in tf.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)
        return lines
    except Exception:
        return []


# ── Single slide → QMD text ───────────────────────────────────────────────────

def render_slide(
    slide,
    slide_num: int,
    lect_num: int,
    images_dir: Path,
    images_rel: str,
    pdf_pages: list[str],
) -> str:
    """Return the QMD text for one slide."""
    lines: list[str] = []

    title = get_slide_title(slide)
    title_safe = title.replace('"', '\\"')

    # ── Heading ──
    if is_section_divider(title):
        lines.append(f"# {title_safe} {{.center .medium-content}}")
    else:
        lines.append(f"# {title_safe or 'Untitled'} {{.medium-content}}")

    # ── Body — iterate all shapes sorted top→bottom, left→right ──
    # Collect and sort first
    all_elems = sorted(
        iter_all_shape_elems(slide),
        key=lambda x: (x[3], x[4]),  # (top, left)
    )

    prev_was_bullet = False

    for tag, elem, ph_type, top, left in all_elems:
        # Skip chrome (title, date, footer, slide number)
        if ph_type in _CHROME_PH_TYPES:
            continue
        # Skip title placeholder (already used for heading)
        if ph_type in (1, 3):
            continue

        # ── Table (graphicFrame containing a:tbl) ──
        if tag == _t(P_NS, "graphicFrame") and _graphicframe_is_table(elem):
            lines.append("")
            lines.append(extract_table_from_elem(elem, slide_num))
            lines.append("")
            prev_was_bullet = False
            continue

        # ── Picture ──
        if tag == _t(P_NS, "pic"):
            # Embedded pictures are handled via the ./Images directory below
            continue

        # ── Text / math shape (p:sp) ──
        if tag == _t(P_NS, "sp"):
            txBody = elem.find(_t(P_NS, "txBody"))
            if txBody is None:
                continue

            is_body = ph_type not in (None,)  # placeholder → body context
            # Non-placeholder (free text box): use body=False so we don't
            # auto-bullet everything; rely only on explicit bullet markers.
            # BUT: if it has no placeholder AND contains only very short text
            # (diagram labels like "q*", "MR"), skip it when the slide has images.
            if ph_type is None:
                # Check if this looks like a diagram label (very short, no paragraphs with >3 words)
                all_text = "".join(
                    p.find(_t(A_NS, "t")).text or ""
                    for p in elem.findall(f".//{_t(A_NS, 't')}")
                    if p.find(_t(A_NS, "t")) is not None  # always true — defensive
                )
                # Simpler: count all text characters
                all_text_chars = sum(
                    len(t.text or "")
                    for t in elem.findall(f".//{_t(A_NS, 't')}")
                )
                has_slide_image = bool(find_slide_images(images_dir, lect_num, slide_num))
                if all_text_chars < 40 and has_slide_image:
                    continue  # skip diagram label on an image slide

            items = extract_text_from_txbody(txBody, is_body=(ph_type is not None))

            for item in items:
                if item.kind == "empty":
                    continue

                if item.kind == "math_display":
                    lines.append("")
                    lines.append("$$")
                    lines.append(item.content)
                    lines.append("$$")
                    lines.append("")
                    prev_was_bullet = False

                elif item.kind == "text":
                    if item.is_bullet or item.level > 0:
                        indent = "  " * item.level
                        if item.list_label:
                            lines.append(f"{indent}{item.list_label} {item.content}")
                        else:
                            lines.append(f"{indent}- {item.content}")
                        prev_was_bullet = True
                    else:
                        if prev_was_bullet:
                            lines.append("")
                        lines.append(item.content)
                        lines.append("")
                        prev_was_bullet = False

    # ── Slide images (matched by Lect{N}_Slide{M}) ──
    slide_imgs = find_slide_images(images_dir, lect_num, slide_num)
    for img_path in slide_imgs:
        rel = f"{images_rel}/{img_path.name}"

        # Caption: concise fragment derived from slide title
        caption = f"{title} — figure" if title else f"Slide {slide_num} figure"

        lines.append("")
        lines.append(f'![{caption}]({rel}){{fig-alt="TODO"}}')
        lines.append("")

    # ── Speaker notes ──
    notes = get_slide_notes(slide)
    if notes:
        lines.append("")
        lines.append("::: {.notes}")
        for note in notes:
            lines.append(f"- {note}")
        lines.append(":::")

    return "\n".join(lines)


# ── File discovery helpers ────────────────────────────────────────────────────

def get_lect_num(pptx_path: Path) -> int:
    """Extract lecture number from filename ('Class 10 - ...' → 10)."""
    m = re.search(r"[Cc]lass\s+(\d+)", pptx_path.stem)
    return int(m.group(1)) if m else 0


def find_matching_pdf(pptx_path: Path) -> Path | None:
    """Locate a handout PDF adjacent to the PPTX."""
    d = pptx_path.parent
    stem = pptx_path.stem
    # Strip trailing " fv" variant common in this project
    stem_no_fv = re.sub(r"\s+fv$", "", stem, flags=re.IGNORECASE)

    candidates = [
        d / f"{stem} handout.pdf",
        d / f"{stem_no_fv} handout.pdf",
        d / f"{stem}.pdf",
    ]
    for c in candidates:
        if c.exists():
            return c

    # Fuzzy: any PDF sharing the "Class N" prefix
    m = re.search(r"[Cc]lass\s+(\d+)", stem)
    if m:
        prefix = f"Class {m.group(1)} "
        for f in sorted(d.glob("*.pdf")):
            if f.stem.lower().startswith(prefix.lower()):
                return f
    return None


def find_images_dir(pptx_path: Path) -> Path:
    """Return the Images directory, trying both capitalizations."""
    for name in ("Images", "images"):
        p = pptx_path.parent / name
        if p.exists():
            return p
    return pptx_path.parent / "Images"  # default (may not exist)


# ── Master conversion ─────────────────────────────────────────────────────────

def convert(
    pptx_path: Path,
    pdf_path: Path | None = None,
    images_dir: Path | None = None,
) -> str:
    """Convert the PPTX and return the full QMD string."""
    basename = pptx_path.stem
    lect_num = get_lect_num(pptx_path)

    if images_dir is None:
        images_dir = find_images_dir(pptx_path)

    # Relative path for image references inside the QMD
    images_rel = f"./{images_dir.name}"

    if pdf_path is None:
        pdf_path = find_matching_pdf(pptx_path)
    pdf_pages = extract_pdf_pages(pdf_path) if pdf_path else []

    if pdf_path:
        print(f"  PDF        : {pdf_path.name}", file=sys.stderr)
    if images_dir.exists():
        img_count = sum(1 for f in images_dir.iterdir() if f.is_file())
        print(f"  Images dir : {images_dir} ({img_count} files)", file=sys.stderr)
    else:
        print(f"  Images dir : {images_dir} (not found)", file=sys.stderr)

    prs = Presentation(str(pptx_path))

    # ── Detect cover slide (ctrTitle placeholder, ph_type == 3) ──
    # If the first PPTX slide is a cover/title slide, use its ctrTitle text
    # as the QMD title and skip it from the body slides.
    slides = list(prs.slides)
    first_slide_is_cover = False
    title_text = basename  # default: filename stem

    if slides:
        first_slide = slides[0]
        for tag, elem, ph_type, *_ in iter_all_shape_elems(first_slide):
            if tag == _t(P_NS, "sp") and ph_type == 3:  # ctrTitle
                txBody = elem.find(_t(P_NS, "txBody"))
                if txBody is not None:
                    texts = []
                    for p in txBody.findall(_t(A_NS, "p")):
                        toks = _extract_para_tokens(p)
                        texts.append(_tokens_to_md(toks))
                    t = _normalize_title(" ".join(texts))
                    if t:
                        title_text = t.title()
                        first_slide_is_cover = True
                break

    # ── Preamble ──
    title_safe = title_text.replace('"', '\\"')
    preamble = (
        "---\n"
        f'title: "{title_safe}"\n'
        'subtitle: "Economics 100B - Intermediate Microeconomics II"\n'
        'author: "Prof. Muehlegger"\n'
        "format:\n"
        "  revealjs:\n"
        "    slide-level: 1\n"
        "---"
    )

    parts = [preamble]

    start_idx = 1 if first_slide_is_cover else 0
    for i, slide in enumerate(slides[start_idx:], start=start_idx + 1):
        slide_text = render_slide(
            slide, i, lect_num, images_dir, images_rel, pdf_pages
        )
        parts.append(slide_text)

    return "\n\n".join(parts) + "\n"


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert a PPTX lecture file to Quarto Markdown (.qmd)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        "pptx_positional", nargs="?", metavar="PPTX",
        help="Path to the .pptx file (positional form)",
    )
    parser.add_argument(
        "--pptx", metavar="PPTX",
        help="Path to the .pptx file (named form)",
    )
    parser.add_argument(
        "--pdf", metavar="PDF",
        help="Path to the handout PDF (default: auto-detected next to PPTX)",
    )
    parser.add_argument(
        "--images", metavar="DIR",
        help="Path to the images directory (default: ./Images next to PPTX)",
    )
    parser.add_argument(
        "--out", metavar="QMD",
        help="Output .qmd path (default: same dir and stem as PPTX)",
    )

    args = parser.parse_args()

    pptx_input = args.pptx_positional or args.pptx
    if not pptx_input:
        parser.error("Provide the path to a .pptx file as a positional or --pptx argument.")

    pptx_path = Path(pptx_input).resolve()
    if not pptx_path.exists():
        sys.exit(f"File not found: {pptx_path}")
    if pptx_path.suffix.lower() != ".pptx":
        sys.exit(f"Expected a .pptx file, got: {pptx_path.name}")

    pdf_path    = Path(args.pdf).resolve()    if args.pdf    else None
    images_dir  = Path(args.images).resolve() if args.images else None
    out_path    = Path(args.out)              if args.out    else pptx_path.with_suffix(".qmd")

    print(f"Converting : {pptx_path.name}", file=sys.stderr)
    print(f"Output     : {out_path}",        file=sys.stderr)

    qmd = convert(pptx_path, pdf_path, images_dir)
    out_path.write_text(qmd, encoding="utf-8")
    print(f"Done → {out_path}", file=sys.stderr)


if __name__ == "__main__":
    main()
